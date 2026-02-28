import sys
from pptx import Presentation

file_path = r"C:\Users\Admin\Downloads\Project Submission Deck _ AMD Slingshot.pptx"

with open("ppt_output.txt", "w", encoding="utf-8") as f:
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            f.write(f"\n===== Slide {i+1} =====\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    f.write(shape.text.strip() + "\n")
    except Exception as e:
        f.write(f"Error reading pptx: {e}\n")
